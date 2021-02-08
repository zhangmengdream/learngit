from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# https://www.jianshu.com/p/d44ad7a9d3aa
def pptx_data_parse(item):
    prs = Presentation(item['local_pdf_url'])
    txt_oa = []
    for x in range(len(prs.slides)):
        for shape in prs.slides[x].shapes:
            if hasattr(shape, "text"):
                row_text = shape.text.encode('utf-8').strip().decode()
                txt_oa.extend(row_text)
        group_shapes = [shp for shp in prs.slides[x].shapes
                        if shp.shape_type == MSO_SHAPE_TYPE.GROUP]
        for group_shape in group_shapes:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    row_text = shape.text.encode('utf-8').strip().decode()
                    txt_oa.extend(row_text)
    txt_oa = ''.join(txt_oa).strip()
    return txt_oa


if __name__ == '__main__':
    item = {"local_pdf_url": r'/home/mikefiles/scrapy_project/5f94410813826de238a87ddf2f1e21cf.pptx'}
    text_list = pptx_data_parse(item)
    print(text_list)
